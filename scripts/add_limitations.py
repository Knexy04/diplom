#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Дописывает концептуальные пункты в слайд «Ограничения и направления развития».
Существующие пункты и прочие шейпы не затрагиваются."""
import os
import shutil
import datetime
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn

PPTX = "/Users/nikitamarkin/diplom/итог/Презентация_Маркин_15.pptx"
TS = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")

LIM_NEW = [
    "•  близость ≠ сопровождение: критерий проверяет геометрию, а не реальную связь (посторонний взрослый рядом может ложно снять тревогу)",
    "•  жёсткий порог возраста 12 лет: подростки на границе классифицируются неустойчиво",
    "•  радиус сопровождения задан в пикселях и не учитывает перспективу камеры",
]
DIR_NEW = [
    "•  моделирование устойчивых пар «ребёнок–взрослый» во времени (кто реально ведёт ребёнка)",
    "•  калибровка плоскости пола (гомография): радиус сопровождения в метрах",
    "•  регрессия возраста вместо бинарного порога — мягкая оценка с уверенностью",
]


def append_bullets(tf, items):
    txBody = tf._txBody
    ps = txBody.findall(qn('a:p'))
    template = ps[-1]
    for text in items:
        new_p = deepcopy(template)
        rs = new_p.findall(qn('a:r'))
        for extra in rs[1:]:
            new_p.remove(extra)
        t = rs[0].find(qn('a:t'))
        t.text = text
        txBody.append(new_p)


shutil.copy(PPTX, PPTX + f".bak_lim_{TS}")
prs = Presentation(PPTX)

slide = None
for s in prs.slides:
    for sh in s.shapes:
        if sh.has_text_frame and "Ограничения и направления развития" in sh.text_frame.text:
            slide = s
            break
    if slide:
        break
if slide is None:
    raise RuntimeError("Слайд не найден")

lim_box = dir_box = None
for sh in slide.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text
    if "классификатора лица" in t and t.strip().startswith("•"):
        lim_box = sh
    elif "балансировка выборок" in t and t.strip().startswith("•"):
        dir_box = sh
if lim_box is None or dir_box is None:
    raise RuntimeError(f"Колонки не найдены: lim={lim_box}, dir={dir_box}")

append_bullets(lim_box.text_frame, LIM_NEW)
append_bullets(dir_box.text_frame, DIR_NEW)
# увеличить высоту блоков под новый объём текста
for box in (lim_box, dir_box):
    box.height = Inches(4.6)

prs.save(PPTX)
print("Пункты добавлены. Бэкап:", PPTX + f".bak_lim_{TS}")
